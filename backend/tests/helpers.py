"""测试共享辅助：用户注册、伪造文件字节、任务轮询等待、假 embedding。"""

import io
import time

# 最小合法文件签名（仅满足 magic bytes 校验，非真实媒体文件）
FAKE_MP4 = b"\x00\x00\x00\x18ftypisom\x00\x00\x00\x00" + b"\x00" * 64
FAKE_WAV = b"RIFF\x24\x00\x00\x00WAVEfmt " + b"\x00" * 64
FAKE_MP3 = b"ID3\x04\x00\x00\x00\x00\x00\x00" + b"\x00" * 64
FAKE_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d4948445200000001000000010806000000"
    "1f15c4890000000d4944415478da63fcffff3f030005fe02fea72d993d"
    "0000000049454e44ae426082"
)


def register(client, email="tester@example.com", password="password123", full_name="张三"):
    resp = client.post(
        "/api/auth/register",
        json={"email": email, "password": password, "full_name": full_name},
    )
    assert resp.status_code == 200, resp.text
    return resp.json()


def auth_headers(tokens) -> dict:
    return {"Authorization": f"Bearer {tokens['access_token']}"}


def generate_inspection(client, tokens, remarks=None):
    data = {"remarks": remarks} if remarks else {}
    resp = client.post(
        "/api/inspection-record/generate",
        headers=auth_headers(tokens),
        files={"video": ("scene.mp4", FAKE_MP4, "video/mp4")},
        data=data,
    )
    assert resp.status_code == 200, resp.text
    return resp.json()["task_id"]


def wait_task(client, tokens, task_id, timeout=15.0) -> dict:
    """轮询任务直至终态（completed/failed/cancelled）。"""
    deadline = time.monotonic() + timeout
    while time.monotonic() < deadline:
        resp = client.get(f"/api/tasks/{task_id}", headers=auth_headers(tokens))
        assert resp.status_code == 200, resp.text
        body = resp.json()
        if body["status"] in ("completed", "failed", "cancelled"):
            return body
        time.sleep(0.05)
    raise AssertionError(f"任务 {task_id} 在 {timeout}s 内未到达终态")


def make_role(user_id: str, role: str) -> None:
    """直接修改用户角色（测试辅助，绕过管理端点）。"""
    import uuid

    from app.db import SessionLocal
    from app.models.user import User

    session = SessionLocal()
    try:
        user = session.get(User, uuid.UUID(user_id))
        user.role = role
        session.commit()
    finally:
        session.close()


def make_admin(user_id: str) -> None:
    """直接把用户角色改为 admin。"""
    make_role(user_id, "admin")


def make_test_video_bytes(seconds: int = 5, with_audio: bool = True) -> bytes:
    """用 imageio-ffmpeg 自带的 ffmpeg 合成真实短视频（testsrc + 可选正弦音轨）。

    供管线端到端测试使用：文件满足扩展名/MIME/签名校验且 ffmpeg 可真实抽帧。
    """
    import subprocess
    import tempfile
    from pathlib import Path

    import imageio_ffmpeg

    exe = imageio_ffmpeg.get_ffmpeg_exe()
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "scene.mp4"
        args = [
            exe, "-y", "-hide_banner", "-loglevel", "error",
            "-f", "lavfi", "-i", f"testsrc=duration={seconds}:size=64x64:rate=10",
        ]
        if with_audio:
            args += ["-f", "lavfi", "-i", f"sine=frequency=440:duration={seconds}"]
        args += ["-shortest", str(out)]
        proc = subprocess.run(args, capture_output=True)
        assert proc.returncode == 0, proc.stderr.decode(errors="replace")
        return out.read_bytes()


def make_test_wav_bytes(seconds: int = 3) -> bytes:
    """用 ffmpeg 合成真实 WAV（正弦音），供询问记录管线测试使用。"""
    import subprocess
    import tempfile
    from pathlib import Path

    import imageio_ffmpeg

    exe = imageio_ffmpeg.get_ffmpeg_exe()
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "audio.wav"
        proc = subprocess.run(
            [exe, "-y", "-hide_banner", "-loglevel", "error",
             "-f", "lavfi", "-i", f"sine=frequency=440:duration={seconds}", str(out)],
            capture_output=True,
        )
        assert proc.returncode == 0, proc.stderr.decode(errors="replace")
        return out.read_bytes()


def fake_embed(texts: list[str]) -> list[list[float]]:
    """确定性假 embedding：按字符分布构造 16 维向量（测试专用，非真实模型）。"""
    vectors = []
    for text in texts:
        vec = [0.0] * 16
        for ch in text:
            vec[ord(ch) % 16] += 1.0
        vectors.append(vec)
    return vectors


def make_minimal_pdf(text: str) -> bytes:
    """构造含文本的最小合法 PDF（ASCII 文本，pypdf 可抽取）。"""
    stream = f"BT /F1 24 Tf 100 700 Td ({text}) Tj ET".encode()
    out = b"%PDF-1.4\n"
    offsets = []
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    pos = len(out)
    for i, obj in enumerate(objs, 1):
        offsets.append(pos)
        entry = f"{i} 0 obj\n".encode() + obj + b"\nendobj\n"
        out += entry
        pos += len(entry)
    xref = b"xref\n0 6\n0000000000 65535 f \n"
    for off in offsets:
        xref += f"{off:010d} 00000 n \n".encode()
    trailer = (
        b"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n"
        + str(pos).encode()
        + b"\n%%EOF"
    )
    return out + xref + trailer


def make_docx(paragraphs: list[str]) -> bytes:
    import docx

    document = docx.Document()
    for p in paragraphs:
        document.add_paragraph(p)
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def make_pptx(slides: list[str]) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    for text in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = text
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()
