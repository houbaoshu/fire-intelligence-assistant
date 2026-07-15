from __future__ import annotations

import shutil
import subprocess
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.core.exceptions import AppError


def parse_document(path: Path, extension: str) -> list[tuple[str, dict[str, Any]]]:
    extension = extension.lower()
    try:
        if extension == ".pdf":
            reader = PdfReader(path)
            return [
                (page.extract_text() or "", {"page": index})
                for index, page in enumerate(reader.pages, 1)
                if (page.extract_text() or "").strip()
            ]
        if extension == ".docx":
            document = Document(str(path))
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
            for table in document.tables:
                text += "\n" + "\n".join(
                    "\t".join(cell.text for cell in row.cells) for row in table.rows
                )
            return [(text, {"section": "document"})]
        if extension == ".pptx":
            presentation = Presentation(str(path))
            results = []
            for index, slide in enumerate(presentation.slides, 1):
                text = "\n".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text
                )
                if text.strip():
                    results.append((text, {"page": index}))
            return results
        if extension in {".doc", ".ppt"}:
            return _parse_legacy(path, extension)
    except AppError:
        raise
    except Exception as error:
        raise AppError(
            status_code=422,
            code="DOCUMENT_PARSE_FAILED",
            message="The document could not be parsed. It may be encrypted or corrupted.",
        ) from error
    raise AppError(
        status_code=415,
        code="UNSUPPORTED_DOCUMENT_TYPE",
        message="The document type is not supported by the parser.",
    )


def _parse_legacy(path: Path, extension: str) -> list[tuple[str, dict[str, Any]]]:
    if extension == ".doc" and (antiword := shutil.which("antiword")):
        result = subprocess.run(
            [antiword, str(path)], capture_output=True, text=True, timeout=60, check=False
        )
        if result.returncode == 0 and result.stdout.strip():
            return [(result.stdout, {"section": "document"})]
    if office := (shutil.which("libreoffice") or shutil.which("soffice")):
        target_extension = ".docx" if extension == ".doc" else ".pptx"
        with TemporaryDirectory(prefix="fip-legacy-") as temporary:
            result = subprocess.run(
                [
                    office,
                    "--headless",
                    "--convert-to",
                    target_extension.lstrip("."),
                    "--outdir",
                    temporary,
                    str(path),
                ],
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
            converted = Path(temporary) / f"{path.stem}{target_extension}"
            if result.returncode == 0 and converted.is_file():
                return parse_document(converted, target_extension)
            raise AppError(
                status_code=422,
                code="LEGACY_CONVERSION_FAILED",
                message=f"The legacy {extension} document could not be converted safely.",
            )
    raise AppError(
        status_code=503,
        code="LEGACY_PARSER_UNAVAILABLE",
        message=(
            f"A legacy {extension} converter is not installed. Convert the file to DOCX or PPTX."
        ),
    )
