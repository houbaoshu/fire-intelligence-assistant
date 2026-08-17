"""PPTX 解析（python-pptx）：每页幻灯片一个文本块，页码为幻灯片序号。"""

import io

from app.rag.parsers.base import ParsedBlock, ParsedDocument, parse_error


def parse_pptx(data: bytes) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError:
        raise parse_error("PPTX 解析组件未安装，请联系管理员")
    try:
        presentation = Presentation(io.BytesIO(data))
        blocks = []
        for i, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text)
            if texts:
                blocks.append(ParsedBlock(text="\n".join(texts), page=i))
    except Exception as exc:
        from app.core.exceptions import AppException

        if isinstance(exc, AppException):
            raise
        raise parse_error("PPTX 文件损坏或格式不受支持，无法解析")
    return ParsedDocument(blocks=blocks)
