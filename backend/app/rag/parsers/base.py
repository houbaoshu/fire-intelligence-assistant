"""Document parser for extracting text from various file formats.

Supports PDF, DOCX, PPTX, and plain text files. Binary formats such as
legacy DOC and PPT are rejected with a clear error message.
"""

from __future__ import annotations

from pathlib import Path

from app.core.exceptions import AppException
from app.core.logging import get_logger

logger = get_logger(__name__)

# Mapping of file extensions to parser method names
_EXTENSION_MAP: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt": "text",
    ".md": "text",
    ".doc": "legacy_doc",
    ".ppt": "legacy_ppt",
}

# MIME type fallback mapping
_MIME_MAP: dict[str, str] = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/plain": "text",
    "text/markdown": "text",
    "application/msword": "legacy_doc",
    "application/vnd.ms-powerpoint": "legacy_ppt",
}


class DocumentParser:
    """Parse documents into text content.

    Supports PDF, DOCX, PPTX, TXT, and MD files.  Legacy binary formats
    (DOC, PPT) are not supported and will raise an error suggesting
    conversion first.
    """

    async def parse(self, file_path: str, mime_type: str | None = None) -> str:
        """Parse a document file into plain text.

        The file type is detected from the extension first.  If the
        extension is unrecognised, *mime_type* is used as a fallback.

        Args:
            file_path: Absolute path to the document file.
            mime_type: Optional MIME type string used as fallback.

        Returns:
            Extracted plain text content.

        Raises:
            AppException: If the file type is unsupported or the file
                cannot be parsed.
        """
        path = Path(file_path)

        if not path.exists():
            raise AppException(
                f"File not found: {file_path}",
                details={"code": "FILE_NOT_FOUND"},
            )

        # Determine parser key from extension, then MIME fallback
        ext = path.suffix.lower()
        parser_key = _EXTENSION_MAP.get(ext)

        if parser_key is None and mime_type:
            parser_key = _MIME_MAP.get(mime_type)

        if parser_key is None:
            raise AppException(
                f"Unsupported file type: {ext or mime_type}",
                details={
                    "code": "UNSUPPORTED_FILE_TYPE",
                    "supported_types": [".pdf", ".docx", ".pptx", ".txt", ".md"],
                },
            )

        logger.info(
            "Parsing document",
            extra={"file_path": file_path, "type": parser_key},
        )

        try:
            if parser_key == "pdf":
                return self._parse_pdf(file_path)
            elif parser_key == "docx":
                return self._parse_docx(file_path)
            elif parser_key == "pptx":
                return self._parse_pptx(file_path)
            elif parser_key == "text":
                return self._parse_text(file_path)
            elif parser_key == "legacy_doc":
                raise AppException(
                    "Legacy .doc format is not supported. Please convert to .docx first.",
                    details={"code": "UNSUPPORTED_FILE_TYPE"},
                )
            elif parser_key == "legacy_ppt":
                raise AppException(
                    "Legacy .ppt format is not supported. Please convert to .pptx first.",
                    details={"code": "UNSUPPORTED_FILE_TYPE"},
                )
            else:
                raise AppException(
                    f"No parser implemented for type: {parser_key}",
                    details={"code": "PARSER_NOT_FOUND"},
                )
        except AppException:
            raise
        except Exception as exc:
            logger.error("Failed to parse document: %s", exc)
            raise AppException(
                f"Failed to parse document: {path.name}",
                details={"code": "PARSE_ERROR", "original_error": str(exc)},
            ) from exc

    def _parse_pdf(self, file_path: str) -> str:
        """Extract text from a PDF file using PyPDF2.

        Falls back to raw binary text extraction if PyPDF2 is not
        available.
        """
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(file_path)
            pages: list[str] = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)

            result = "\n\n".join(pages)
            if not result.strip():
                logger.warning("PDF parsed but no text extracted: %s", file_path)
            return result

        except ImportError:
            logger.warning("PyPDF2 not installed, falling back to raw text extraction")
            return self._fallback_text_extract(file_path)

    def _parse_docx(self, file_path: str) -> str:
        """Extract text from a DOCX file using python-docx."""
        try:
            from docx import Document
        except ImportError as exc:
            raise AppException(
                "python-docx is required for parsing .docx files. "
                "Install it with: pip install python-docx",
                details={"code": "DEPENDENCY_MISSING"},
            ) from exc

        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    def _parse_pptx(self, file_path: str) -> str:
        """Extract text from a PPTX file using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise AppException(
                "python-pptx is required for parsing .pptx files. "
                "Install it with: pip install python-pptx",
                details={"code": "DEPENDENCY_MISSING"},
            ) from exc

        prs = Presentation(file_path)
        slides_text: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            shapes_text: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            shapes_text.append(text)
            if shapes_text:
                slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(shapes_text))

        return "\n\n".join(slides_text)

    def _parse_text(self, file_path: str) -> str:
        """Read a plain text or markdown file."""
        encodings = ("utf-8", "gbk", "gb2312", "latin-1")
        for encoding in encodings:
            try:
                with open(file_path, encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue

        raise AppException(
            f"Unable to decode text file: {file_path}",
            details={"code": "ENCODING_ERROR"},
        )

    @staticmethod
    def _fallback_text_extract(file_path: str) -> str:
        """Best-effort text extraction from a binary file.

        Used only when PyPDF2 is not available.  Reads raw bytes and
        decodes what it can.
        """
        try:
            with open(file_path, "rb") as f:
                raw = f.read()
            # Attempt UTF-8 decode, ignoring errors
            return raw.decode("utf-8", errors="ignore")
        except Exception:
            return ""
